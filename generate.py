from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Define slide content
slides_content = [
    ("Real-Time Communication with Server-Sent Events (SSE)", "Enabling one-way real-time updates from server to client\nPresented by: [Your Name]\nDate: [Presentation Date]"),
    ("Agenda", "1. What are Server-Sent Events?\n2. How SSE Works\n3. Benefits of SSE\n4. SSE vs WebSockets vs Polling\n5. Ideal Use Cases\n6. Implementation Overview\n7. Demo / Architecture Diagram\n8. Limitations & Considerations\n9. Q&A"),
    ("What Are Server-Sent Events?", "- A W3C standard for pushing real-time updates from server to browser\n- Built on HTTP (not a new protocol)\n- One-way communication: Server ➝ Client\n- Lightweight alternative to WebSockets for specific use cases"),
    ("How It Works", "1. Client opens an EventSource connection to the server.\n2. Server keeps the HTTP connection open and streams updates.\n3. Events are formatted as text/event-stream.\n4. Automatically handles reconnection on client side."),
    ("Key Benefits", "- ✅ Built-in browser support (no extra libraries for client)\n- ✅ Auto-reconnect, last event ID for resume support\n- ✅ Efficient for unidirectional real-time updates\n- ✅ Uses HTTP/1.1 or HTTP/2, integrates with proxies and firewalls"),
    ("SSE vs Alternatives", "Feature | SSE | WebSocket | Polling\n--------|-----|----------|--------\nDirection | Server ➝ Client | Bidirectional | Client-initiated\nProtocol | HTTP | Custom (ws://) | HTTP\nComplexity | Low | Medium | Low\nReconnect Support | Built-in | Manual | Not needed\nUse Case | Notifications, Feeds | Chat, Gaming | Low-frequency updates"),
    ("Ideal Use Cases", "- ✅ Live notifications (e.g., alerts, system messages)\n- ✅ Stock price or sensor updates\n- ✅ Monitoring dashboards\n- ✅ Streaming logs/events\n- ⚠️ Not for chat or heavy bi-directional apps (use WebSockets)"),
    ("Architecture Overview", "Client (EventSource) → HTTP → Backend SSE Endpoint\n↳ Generates events\n↳ Streams via text/event-stream"),
    ("Sample Backend (Spring Boot)", "@GetMapping(\"/events\")\npublic SseEmitter streamEvents() {\n    SseEmitter emitter = new SseEmitter();\n    executor.submit(() -> {\n        try {\n            while (true) {\n                emitter.send(SseEmitter.event().data(\"Heartbeat\"));\n                Thread.sleep(1000);\n            }\n        } catch (Exception ex) {\n            emitter.completeWithError(ex);\n        }\n    });\n    return emitter;\n}"),
    ("Limitations & Considerations", "- ❌ One-way only (use WebSockets if you need two-way)\n- ❌ No built-in binary support\n- ❌ Mobile support may vary\n- ❗ Connection limits on some browsers (~6 per domain)\n- 🔒 Requires heartbeat or keep-alive to avoid idle disconnects"),
    ("Summary", "- SSE is simple, efficient, and reliable for real-time server ➝ client updates\n- Best suited for lightweight, uni-directional use cases\n- Easier to implement and manage compared to WebSockets\n- Consider use case carefully when choosing between SSE, WS, or polling"),
    ("Q&A", "💬 “Any questions on how SSE can benefit our real-time architecture?”")
]

# Add slides to presentation
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

# Save the presentation
pptx_file = "/mnt/data/Server_Sent_Events_Presentation.pptx"
prs.save(pptx_file)
